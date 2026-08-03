"""
Build a branded pptx reference template (BTS colors, pulled from
agent-charles-dashboard.html's --color-primary/--color-accent) for pandoc's
`--reference-doc` flag. Pandoc reuses whichever slide layout matches by name
(Title Slide / Title and Content / Section Header / Title Only) for every
`# heading` slide it generates, so branding those four layouts here is
enough to brand every deck built from a slides.md outline.

Usage:
    python build_pptx_template.py
Produces:
    templates/bts-reference.pptx
"""

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

NAVY = RGBColor(0x1E, 0x3A, 0x5F)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GRAY_DARK = RGBColor(0x37, 0x41, 0x51)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

OUT_PATH = Path("templates/bts-reference.pptx")


def solid_bg(layout, color):
    fill = layout.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def style_placeholder(layout, ph_idx, size, color, bold=True):
    ph = layout.placeholders[ph_idx]

    # Sets the prompt text ("Click to edit...") shown when editing the layout
    # itself — cosmetic only, real slide content does NOT inherit from this.
    p = ph.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color

    # The actual inheritance path for text typed into slides built from this
    # layout is <a:lstStyle><a:lvl1pPr><a:defRPr>, not the paragraph's own
    # pPr/defRPr above — without this, PowerPoint falls back past the (empty)
    # lstStyle straight to the slide master's default titleStyle/bodyStyle.
    # Body/text placeholders (unlike title placeholders) ship with a
    # pre-populated lstStyle (bullet levels) that already has an lvl1pPr —
    # appending a second one is a duplicate element that corrupts the file
    # (PowerPoint silently refuses to open it), so replace it if present.
    lst_style = ph.text_frame._txBody.find(qn("a:lstStyle"))
    existing_lvl1 = lst_style.find(qn("a:lvl1pPr"))
    if existing_lvl1 is not None:
        lst_style.remove(existing_lvl1)
    lvl1 = etree.Element(qn("a:lvl1pPr"))
    lst_style.insert(0, lvl1)
    def_rpr = etree.SubElement(lvl1, qn("a:defRPr"))
    def_rpr.set("sz", str(int(size * 100)))
    def_rpr.set("b", "1" if bold else "0")
    fill = etree.SubElement(def_rpr, qn("a:solidFill"))
    clr = etree.SubElement(fill, qn("a:srgbClr"))
    clr.set("val", str(color))


# python-pptx's LayoutShapes has no add_shape/add_textbox (only SlideShapes
# does) — build the shape on a throwaway slide, then graft its XML element
# onto the target layout's shape tree. Pandoc's --reference-doc only reads
# masters/layouts and discards whatever slides already exist in the file, so
# the scratch slide left behind in the saved template is harmless.
def _scratch_shapes(prs):
    return prs.slides.add_slide(prs.slide_layouts[6]).shapes


def _graft(shape, layout):
    # A shape built on the scratch slide gets shape IDs scoped to that slide's
    # own tree (starting back at 2), which collide with the layout's existing
    # placeholder IDs once moved over — duplicate p:cNvPr ids are invalid
    # OOXML and made PowerPoint refuse to open the file. Renumber on the way in.
    tree = layout.shapes.element
    existing_ids = [int(el.get("id")) for el in tree.iter(qn("p:cNvPr"))]
    new_id = max(existing_ids, default=0) + 1
    shape._element.find(".//" + qn("p:cNvPr")).set("id", str(new_id))
    tree.append(shape._element)


def wordmark(prs, layout, color, slide_w, slide_h):
    width, height = Emu(1500000), Emu(300000)
    box = _scratch_shapes(prs).add_textbox(
        slide_w - width - Emu(228600), slide_h - height - Emu(160000), width, height
    )
    tf = box.text_frame
    tf.paragraphs[0].text = "BTS"
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = color
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _graft(box, layout)


def accent_bar(prs, layout, top, slide_w, height=Emu(45720)):
    bar = _scratch_shapes(prs).add_shape(MSO_SHAPE.RECTANGLE, Emu(0), top, slide_w, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()
    _graft(bar, layout)


def main() -> None:
    prs = Presentation()
    w, h = prs.slide_width, prs.slide_height

    title_slide = prs.slide_layouts[0]
    subtitle_ph = title_slide.placeholders[1]
    solid_bg(title_slide, NAVY)
    style_placeholder(title_slide, 0, 40, WHITE)
    style_placeholder(title_slide, 1, 18, AMBER, bold=False)
    accent_bar(prs, title_slide, subtitle_ph.top + subtitle_ph.height, w)
    wordmark(prs, title_slide, AMBER, w, h)

    content = prs.slide_layouts[1]
    title_ph = content.placeholders[0]
    solid_bg(content, WHITE)
    style_placeholder(content, 0, 28, NAVY)
    style_placeholder(content, 1, 18, GRAY_DARK, bold=False)
    accent_bar(prs, content, title_ph.top + title_ph.height, w)
    wordmark(prs, content, NAVY, w, h)

    section = prs.slide_layouts[2]
    solid_bg(section, NAVY)
    style_placeholder(section, 0, 34, WHITE)
    style_placeholder(section, 1, 18, AMBER, bold=False)
    wordmark(prs, section, AMBER, w, h)

    title_only = prs.slide_layouts[5]
    title_only_ph = title_only.placeholders[0]
    solid_bg(title_only, WHITE)
    style_placeholder(title_only, 0, 28, NAVY)
    accent_bar(prs, title_only, title_only_ph.top + title_only_ph.height, w)
    wordmark(prs, title_only, NAVY, w, h)

    # The scratch (Blank-layout) slides used to fabricate shapes for grafting
    # are left in place — pandoc's --reference-doc only reads masters/layouts
    # and discards whatever slides already exist, and removing them via the
    # sldIdLst XML directly (without touching presentation.xml.rels/parts)
    # left orphaned slide parts that collided with real python-pptx-added
    # slides on the next save (duplicate slideN.xml part names).

    OUT_PATH.parent.mkdir(exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
